"""
Report Generator
----------------
Produces professional PDF and PPTX validation reports from model data.

PDF:  reportlab Platypus — full structured report for compliance files
PPTX: python-pptx       — executive summary deck for client presentations

Valtior brand colours:
  Navy  #0B1F3A  (RGB 11, 31, 58)
  Teal  #007B7B  (RGB 0, 123, 123)
  Light #F4F7FA
"""

import io
from datetime import datetime

# ── Severity colours ──────────────────────────────────────────────────────────
SEV_COLOURS = {
    "Critical":    (220, 38, 38),    # red-600
    "High":        (234, 88, 12),    # orange-600
    "Medium":      (37, 99, 235),    # blue-600
    "Low":         (22, 163, 74),    # green-600
    "Observation": (107, 114, 128),  # gray-500
}

STATUS_COLOURS = {
    "MET":     (22, 163, 74),     # green
    "PARTIAL": (234, 88, 12),     # orange
    "NOT_MET": (220, 38, 38),     # red
    "UNCLEAR": (107, 114, 128),   # gray
}

STATUS_LABELS = {
    "MET":     "✓ Met",
    "PARTIAL": "~ Partial",
    "NOT_MET": "✗ Not Met",
    "UNCLEAR": "? Unclear",
}

NAVY  = (11, 31, 58)
TEAL  = (0, 123, 123)
LIGHT = (244, 247, 250)
WHITE = (255, 255, 255)


# ─────────────────────────────────────────────────────────────────────────────
# PDF Report
# ─────────────────────────────────────────────────────────────────────────────

def generate_pdf_report(project, model, findings, compliance_results, compliance_summary_data) -> bytes:
    """
    Generate a professional PDF validation report.
    Returns raw PDF bytes.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm, mm
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, HRFlowable, KeepTogether
    )
    from reportlab.lib.colors import Color, HexColor

    buffer = io.BytesIO()
    page_w, page_h = A4

    def rgb(r, g, b):
        return Color(r/255, g/255, b/255)

    c_navy  = rgb(*NAVY)
    c_teal  = rgb(*TEAL)
    c_light = rgb(*LIGHT)
    c_white = rgb(*WHITE)

    def sev_color(sev):
        return rgb(*SEV_COLOURS.get(sev, (107,114,128)))

    def status_color(st):
        return rgb(*STATUS_COLOURS.get(st, (107,114,128)))

    # ── Page template with header/footer ─────────────────────────────────────
    def on_page(canvas, doc):
        canvas.saveState()
        # Header bar
        canvas.setFillColor(c_navy)
        canvas.rect(0, page_h - 1.2*cm, page_w, 1.2*cm, fill=1, stroke=0)
        canvas.setFillColor(c_white)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(1.5*cm, page_h - 0.85*cm, "VALTIOR — Model Validation Report")
        canvas.setFont("Helvetica", 9)
        canvas.drawRightString(page_w - 1.5*cm, page_h - 0.85*cm, "CONFIDENTIAL")
        # Footer
        canvas.setFillColor(rgb(156,163,175))
        canvas.setFont("Helvetica", 8)
        canvas.drawString(1.5*cm, 0.8*cm, f"{project.institution or project.name}")
        canvas.drawCentredString(page_w/2, 0.8*cm, datetime.now().strftime("%d %B %Y"))
        canvas.drawRightString(page_w - 1.5*cm, 0.8*cm, f"Page {doc.page}")
        canvas.restoreState()

    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=1.8*cm, rightMargin=1.8*cm,
        topMargin=2.2*cm, bottomMargin=2*cm,
        title=f"Validation Report — {project.name}",
        author="Valtior",
    )

    styles = getSampleStyleSheet()

    def style(name, **kw):
        return ParagraphStyle(name, **kw)

    S = {
        "h1": style("H1", fontName="Helvetica-Bold", fontSize=22, textColor=c_navy,
                    spaceAfter=6, leading=26),
        "h2": style("H2", fontName="Helvetica-Bold", fontSize=14, textColor=c_navy,
                    spaceBefore=14, spaceAfter=6, leading=18),
        "h3": style("H3", fontName="Helvetica-Bold", fontSize=11, textColor=c_teal,
                    spaceBefore=8, spaceAfter=4, leading=14),
        "body": style("Body", fontName="Helvetica", fontSize=9.5, leading=14,
                      textColor=colors.HexColor("#374151"), spaceAfter=4),
        "small": style("Small", fontName="Helvetica", fontSize=8.5, leading=12,
                       textColor=colors.HexColor("#6B7280")),
        "label": style("Label", fontName="Helvetica-Bold", fontSize=8,
                       textColor=colors.HexColor("#6B7280"), spaceAfter=2),
        "reg": style("Reg", fontName="Helvetica-Oblique", fontSize=8,
                     textColor=c_teal, leading=11),
        "cover_title": style("CoverTitle", fontName="Helvetica-Bold", fontSize=32,
                             textColor=c_white, leading=38, spaceAfter=8),
        "cover_sub": style("CoverSub", fontName="Helvetica", fontSize=14,
                           textColor=rgb(200,220,240), leading=18),
    }

    story = []

    # ─── Cover page ──────────────────────────────────────────────────────────
    # Navy banner
    story.append(Table(
        [[Paragraph(f"Model Validation Report", S["h1"])]],
        colWidths=[page_w - 3.6*cm],
        style=TableStyle([
            ("BACKGROUND", (0,0), (-1,-1), c_navy),
            ("TEXTCOLOR",  (0,0), (-1,-1), c_white),
            ("TOPPADDING",   (0,0), (-1,-1), 18),
            ("BOTTOMPADDING",(0,0), (-1,-1), 18),
            ("LEFTPADDING",  (0,0), (-1,-1), 16),
        ])
    ))
    story.append(Spacer(1, 0.5*cm))

    # Project details table
    meta_rows = [
        ["Project",     project.name],
        ["Institution", project.institution or "—"],
        ["Model owner", project.owner or "—"],
        ["Domain",      project.domain],
        ["Date",        datetime.now().strftime("%d %B %Y")],
        ["Model version", f"v{model.version}"],
        ["Lifecycle stage", (project.lifecycle_stage or "validering").replace("_", " ").title()],
    ]
    meta_table = Table(
        [[Paragraph(k, S["label"]), Paragraph(v, S["body"])] for k, v in meta_rows],
        colWidths=[5*cm, page_w - 3.6*cm - 5*cm],
        style=TableStyle([
            ("BACKGROUND", (0,0), (0,-1), c_light),
            ("ROWBACKGROUNDS", (0,0), (-1,-1), [c_light, c_white]),
            ("FONTNAME",   (0,0), (0,-1), "Helvetica-Bold"),
            ("FONTSIZE",   (0,0), (-1,-1), 9),
            ("TOPPADDING",    (0,0), (-1,-1), 5),
            ("BOTTOMPADDING", (0,0), (-1,-1), 5),
            ("LEFTPADDING",   (0,0), (-1,-1), 8),
            ("GRID", (0,0), (-1,-1), 0.5, rgb(209,213,219)),
        ])
    )
    story.append(meta_table)
    story.append(Spacer(1, 0.5*cm))

    # RAG status box
    rag = compliance_summary_data["rag"]
    rag_color = {"GREEN": rgb(22,163,74), "AMBER": rgb(234,88,12), "RED": rgb(220,38,38)}[rag]
    story.append(Table(
        [[
            Paragraph(f"Overall Status: {rag}", style("RAG", fontName="Helvetica-Bold",
                fontSize=12, textColor=c_white)),
            Paragraph(compliance_summary_data["rag_label"],
                style("RAGLabel", fontName="Helvetica", fontSize=9, textColor=c_white, leading=12)),
        ]],
        colWidths=[5*cm, page_w - 3.6*cm - 5*cm],
        style=TableStyle([
            ("BACKGROUND", (0,0), (-1,-1), rag_color),
            ("TOPPADDING",    (0,0), (-1,-1), 10),
            ("BOTTOMPADDING", (0,0), (-1,-1), 10),
            ("LEFTPADDING",   (0,0), (-1,-1), 12),
            ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ])
    ))

    story.append(PageBreak())

    # ─── Section 1: Executive Summary ────────────────────────────────────────
    story.append(Paragraph("1. Executive Summary", S["h2"]))
    story.append(HRFlowable(width="100%", thickness=1, color=c_teal, spaceAfter=8))

    counts_by_sev = {s: 0 for s in ["Critical","High","Medium","Low","Observation"]}
    for f in findings:
        sev = (f.get("severity") if isinstance(f,dict) else getattr(f,"severity","Medium"))
        if sev in counts_by_sev:
            counts_by_sev[sev] += 1

    sev_rows = [["Severity", "Count", "Description"]]
    sev_descs = {
        "Critical":    "Immediate remediation required — model cannot be used as-is.",
        "High":        "Must be resolved before validation sign-off.",
        "Medium":      "Should be resolved or formally accepted with rationale.",
        "Low":         "Minor weakness — improvement opportunity.",
        "Observation": "Noted for the validation record.",
    }
    for sev, count in counts_by_sev.items():
        if count > 0:
            sev_rows.append([sev, str(count), sev_descs[sev]])

    sev_style = TableStyle([
        ("BACKGROUND", (0,0), (-1,0), c_navy),
        ("TEXTCOLOR",  (0,0), (-1,0), c_white),
        ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",   (0,0), (-1,-1), 9),
        ("TOPPADDING",    (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
        ("LEFTPADDING",   (0,0), (-1,-1), 8),
        ("GRID", (0,0), (-1,-1), 0.5, rgb(209,213,219)),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [c_white, c_light]),
    ])
    # Colour severity column
    for i, (sev, count) in enumerate(counts_by_sev.items(), 1):
        if count > 0 and sev in SEV_COLOURS:
            sev_style.add("TEXTCOLOR", (0,i), (0,i), rgb(*SEV_COLOURS[sev]))
            sev_style.add("FONTNAME",  (0,i), (0,i), "Helvetica-Bold")

    story.append(Table(
        sev_rows,
        colWidths=[3*cm, 2*cm, page_w - 3.6*cm - 5*cm],
        style=sev_style
    ))
    story.append(Spacer(1, 0.4*cm))

    comp = compliance_summary_data["counts"]
    story.append(Paragraph(
        f"The model was assessed against <b>{compliance_summary_data['total']} Swedish regulatory requirements</b>. "
        f"<b>{comp['MET']}</b> requirements are met, <b>{comp['PARTIAL']}</b> partially met, "
        f"and <b>{comp['NOT_MET']}</b> not met ({compliance_summary_data['met_pct']}% compliance rate).",
        S["body"]
    ))

    story.append(PageBreak())

    # ─── Section 2: Swedish Regulatory Compliance Matrix ─────────────────────
    story.append(Paragraph("2. Swedish Regulatory Compliance Matrix", S["h2"]))
    story.append(HRFlowable(width="100%", thickness=1, color=c_teal, spaceAfter=8))
    story.append(Paragraph(
        "Assessment against mandatory requirements under PTL (2017:630), FFFS 2017:11, and SIMPT guidance.",
        S["body"]
    ))
    story.append(Spacer(1, 0.3*cm))

    comp_rows = [["ID", "Law", "Requirement", "Status"]]
    for r in compliance_results:
        comp_rows.append([
            r["id"],
            r["law"],
            r["short"],
            STATUS_LABELS[r["status"]],
        ])

    comp_table_style = TableStyle([
        ("BACKGROUND", (0,0), (-1,0), c_navy),
        ("TEXTCOLOR",  (0,0), (-1,0), c_white),
        ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",   (0,0), (-1,-1), 8.5),
        ("TOPPADDING",    (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
        ("LEFTPADDING",   (0,0), (-1,-1), 6),
        ("GRID", (0,0), (-1,-1), 0.5, rgb(209,213,219)),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [c_white, c_light]),
        ("FONTNAME", (0,1), (0,-1), "Helvetica-Bold"),
        ("TEXTCOLOR", (0,1), (0,-1), c_teal),
    ])
    for i, r in enumerate(compliance_results, 1):
        sc = status_color(r["status"])
        comp_table_style.add("TEXTCOLOR", (3,i), (3,i), sc)
        comp_table_style.add("FONTNAME",  (3,i), (3,i), "Helvetica-Bold")

    story.append(Table(
        comp_rows,
        colWidths=[1.8*cm, 4.2*cm, 8*cm, 2.5*cm],
        style=comp_table_style
    ))

    story.append(PageBreak())

    # ─── Section 3: Findings by Validation Framework Dimension ──────────────
    story.append(Paragraph("3. Findings — Frank Penny / SIMPT Kontrollramverk", S["h2"]))
    story.append(HRFlowable(width="100%", thickness=1, color=c_teal, spaceAfter=6))
    story.append(Paragraph(
        "Findings are organised by the four validation control framework dimensions "
        "(Frank Penny / SIMPT 2024), matching the structure of the validation questions in FFFS 2017:11 §15.",
        S["body"]
    ))
    story.append(Spacer(1, 0.3*cm))

    DIMS = [
        ("modellutformning",      "Modellutformning",       "Model Design",
         "Har modellen ett tydligt syfte och bygger på lämpliga antaganden?"),
        ("data",                  "Data",                   "Data",
         "Är data korrekt och komplett, och har datamappning dokumenterats?"),
        ("implementation_testning","Implementation & Testning","Implementation & Testing",
         "Är modellen implementerad i enlighet med modellutformning och rutiner?"),
        ("styrning_uppföljning",  "Styrning & Uppföljning", "Governance & Monitoring",
         "Finns det tydliga roller och ansvar under hela modellens livscykel?"),
    ]

    severity_order = {"Critical":0, "High":1, "Medium":2, "Low":3, "Observation":4}

    def _f_val(f, attr):
        return f.get(attr) if isinstance(f, dict) else getattr(f, attr, None)

    # Group findings by dimension
    by_dim = {d[0]: [] for d in DIMS}
    untagged = []
    for f in findings:
        dim = _f_val(f, "framework_dimension")
        if dim and dim in by_dim:
            by_dim[dim].append(f)
        else:
            untagged.append(f)
    for k in by_dim:
        by_dim[k].sort(key=lambda f: severity_order.get(_f_val(f,"severity") or "", 99))

    # Dimension RAG
    def _dim_rag(fs):
        sevs = [_f_val(f,"severity") for f in fs]
        if "Critical" in sevs: return "RED"
        if "High" in sevs:     return "AMBER"
        if fs:                 return "GREEN"
        return "NONE"
    dim_rag_colours = {
        "RED":   rgb(220,38,38),
        "AMBER": rgb(234,88,12),
        "GREEN": rgb(22,163,74),
        "NONE":  rgb(156,163,175),
    }
    dim_rag_labels = {
        "RED":   "✗ Issues found",
        "AMBER": "⚠ Review required",
        "GREEN": "✓ No critical issues",
        "NONE":  "— No findings",
    }

    # Framework overview table
    ov_rows = [["Dimension", "EN", "Status", "Findings"]]
    for dim_key, dim_label, dim_en, _ in DIMS:
        rag = _dim_rag(by_dim[dim_key])
        ov_rows.append([dim_label, dim_en, dim_rag_labels[rag], str(len(by_dim[dim_key]))])
    ov_style = TableStyle([
        ("BACKGROUND", (0,0), (-1,0), c_navy),
        ("TEXTCOLOR",  (0,0), (-1,0), c_white),
        ("FONTNAME",   (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE",   (0,0), (-1,-1), 9),
        ("TOPPADDING",    (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
        ("LEFTPADDING",   (0,0), (-1,-1), 8),
        ("GRID", (0,0), (-1,-1), 0.5, rgb(209,213,219)),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [c_white, c_light]),
    ])
    for i, (dim_key, _, _, _) in enumerate(DIMS, 1):
        rag = _dim_rag(by_dim[dim_key])
        ov_style.add("TEXTCOLOR", (2,i), (2,i), dim_rag_colours[rag])
        ov_style.add("FONTNAME",  (2,i), (2,i), "Helvetica-Bold")
    story.append(Table(
        ov_rows,
        colWidths=[5.5*cm, 5.5*cm, 4*cm, 2*cm],
        style=ov_style
    ))
    story.append(Spacer(1, 0.4*cm))

    # Findings per dimension
    finding_counter = 0
    for dim_key, dim_label, dim_en, dim_q in DIMS:
        dim_findings = by_dim[dim_key]
        rag = _dim_rag(dim_findings)
        rag_c = dim_rag_colours[rag]

        # Dimension header
        story.append(KeepTogether([
            Table(
                [[
                    Paragraph(dim_label,
                        style(f"DH_{dim_key}", fontName="Helvetica-Bold", fontSize=11,
                              textColor=c_white, leading=14)),
                    Paragraph(f"{dim_en} · {dim_rag_labels[rag]}",
                        style(f"DHsub_{dim_key}", fontName="Helvetica", fontSize=8.5,
                              textColor=c_white, leading=12)),
                ]],
                colWidths=[5.5*cm, page_w - 3.6*cm - 5.5*cm],
                style=TableStyle([
                    ("BACKGROUND", (0,0), (-1,-1), rag_c),
                    ("TOPPADDING",    (0,0), (-1,-1), 8),
                    ("BOTTOMPADDING", (0,0), (-1,-1), 8),
                    ("LEFTPADDING",   (0,0), (-1,-1), 12),
                    ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                ])
            ),
            Table(
                [[Paragraph(f'"{dim_q}"',
                    style(f"DQ_{dim_key}", fontName="Helvetica-Oblique", fontSize=8.5,
                          textColor=rgb(75,85,99), leading=12))]],
                colWidths=[page_w - 3.6*cm],
                style=TableStyle([
                    ("BACKGROUND", (0,0), (-1,-1), c_light),
                    ("TOPPADDING",    (0,0), (-1,-1), 6),
                    ("BOTTOMPADDING", (0,0), (-1,-1), 6),
                    ("LEFTPADDING",   (0,0), (-1,-1), 12),
                ])
            ),
            Spacer(1, 0.15*cm),
        ]))

        if not dim_findings:
            story.append(Paragraph("No findings in this dimension.", S["small"]))
            story.append(Spacer(1, 0.3*cm))
            continue

        for f in dim_findings:
            finding_counter += 1
            i = finding_counter
            sev = _f_val(f, "severity") or ""
            title = _f_val(f, "title") or ""
            desc  = _f_val(f, "description") or ""
            rec   = _f_val(f, "recommendation") or ""
            reg   = _f_val(f, "regulatory_reference") or ""
            src   = _f_val(f, "source") or ""
            cat   = _f_val(f, "category") or ""
            ann   = _f_val(f, "annotation_status") or "pending"
            sc    = rgb(*SEV_COLOURS.get(sev, (107,114,128)))

            finding_block = [
                Table(
                    [[
                        Paragraph(sev, style(f"Sev{i}", fontName="Helvetica-Bold", fontSize=8, textColor=c_white)),
                        Paragraph(f"F{i:02d} — {title}",
                                   style(f"FT{i}", fontName="Helvetica-Bold", fontSize=10, textColor=c_white, leading=13)),
                        Paragraph(f"[{ann.upper()}]",
                                   style(f"Ann{i}", fontName="Helvetica", fontSize=7.5, textColor=c_white, alignment=TA_RIGHT)),
                    ]],
                    colWidths=[2.2*cm, 11*cm, 3.3*cm],
                    style=TableStyle([
                        ("BACKGROUND", (0,0), (-1,-1), sc),
                        ("TOPPADDING",    (0,0), (-1,-1), 7),
                        ("BOTTOMPADDING", (0,0), (-1,-1), 7),
                        ("LEFTPADDING",   (0,0), (-1,-1), 8),
                        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                    ])
                ),
                Table(
                    [
                        [Paragraph("Finding", S["label"]),        Paragraph(desc or "—", S["body"])],
                        [Paragraph("Recommendation", S["label"]), Paragraph(rec or "—", S["body"])],
                        [Paragraph("Regulatory ref.", S["label"]), Paragraph(reg or "—", S["reg"])],
                        [Paragraph("Source", S["label"]),
                         Paragraph(f"{src.replace('_',' ').title()} · {cat.replace('_',' ')}", S["small"])],
                    ],
                    colWidths=[3.5*cm, page_w - 3.6*cm - 3.5*cm],
                    style=TableStyle([
                        ("BACKGROUND", (0,0), (0,-1), c_light),
                        ("BACKGROUND", (1,0), (1,-1), c_white),
                        ("TOPPADDING",    (0,0), (-1,-1), 5),
                        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
                        ("LEFTPADDING",   (0,0), (-1,-1), 8),
                        ("GRID", (0,0), (-1,-1), 0.5, rgb(229,231,235)),
                        ("VALIGN", (0,0), (-1,-1), "TOP"),
                    ])
                ),
                Spacer(1, 0.35*cm),
            ]
            story.append(KeepTogether(finding_block))

        story.append(Spacer(1, 0.2*cm))

    # Untagged findings (legacy, no dimension set)
    if untagged:
        story.append(Paragraph("Other Findings", S["h3"]))
        untagged.sort(key=lambda f: severity_order.get(_f_val(f,"severity") or "", 99))
        for f in untagged:
            finding_counter += 1
            i = finding_counter
            sev   = _f_val(f, "severity") or ""
            title = _f_val(f, "title") or ""
            desc  = _f_val(f, "description") or ""
            rec   = _f_val(f, "recommendation") or ""
            reg   = _f_val(f, "regulatory_reference") or ""
            src   = _f_val(f, "source") or ""
            cat   = _f_val(f, "category") or ""
            ann   = _f_val(f, "annotation_status") or "pending"
            sc    = rgb(*SEV_COLOURS.get(sev, (107,114,128)))
            story.append(KeepTogether([
                Table(
                    [[
                        Paragraph(sev, style(f"USev{i}", fontName="Helvetica-Bold", fontSize=8, textColor=c_white)),
                        Paragraph(f"F{i:02d} — {title}",
                                   style(f"UFT{i}", fontName="Helvetica-Bold", fontSize=10, textColor=c_white, leading=13)),
                        Paragraph(f"[{ann.upper()}]",
                                   style(f"UAnn{i}", fontName="Helvetica", fontSize=7.5, textColor=c_white, alignment=TA_RIGHT)),
                    ]],
                    colWidths=[2.2*cm, 11*cm, 3.3*cm],
                    style=TableStyle([
                        ("BACKGROUND", (0,0), (-1,-1), sc),
                        ("TOPPADDING",    (0,0), (-1,-1), 7),
                        ("BOTTOMPADDING", (0,0), (-1,-1), 7),
                        ("LEFTPADDING",   (0,0), (-1,-1), 8),
                        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                    ])
                ),
                Table(
                    [
                        [Paragraph("Finding", S["label"]),        Paragraph(desc or "—", S["body"])],
                        [Paragraph("Recommendation", S["label"]), Paragraph(rec or "—", S["body"])],
                        [Paragraph("Regulatory ref.", S["label"]), Paragraph(reg or "—", S["reg"])],
                    ],
                    colWidths=[3.5*cm, page_w - 3.6*cm - 3.5*cm],
                    style=TableStyle([
                        ("BACKGROUND", (0,0), (0,-1), c_light),
                        ("BACKGROUND", (1,0), (1,-1), c_white),
                        ("TOPPADDING",    (0,0), (-1,-1), 5),
                        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
                        ("LEFTPADDING",   (0,0), (-1,-1), 8),
                        ("GRID", (0,0), (-1,-1), 0.5, rgb(229,231,235)),
                        ("VALIGN", (0,0), (-1,-1), "TOP"),
                    ])
                ),
                Spacer(1, 0.35*cm),
            ]))

    story.append(PageBreak())

    # ─── Section 4: Validation Notes ─────────────────────────────────────────
    story.append(Paragraph("4. Validation Notes & Next Steps", S["h2"]))
    story.append(HRFlowable(width="100%", thickness=1, color=c_teal, spaceAfter=8))

    structured = model.structured or {}
    parsing_notes = (structured.get("metadata") or {}).get("parsing_notes", "")
    parsing_conf  = (structured.get("metadata") or {}).get("parsing_confidence", "")

    if parsing_notes:
        story.append(Paragraph("Parsing Notes", S["h3"]))
        story.append(Paragraph(f"Confidence: {parsing_conf or 'unknown'}", S["small"]))
        story.append(Paragraph(parsing_notes, S["body"]))
        story.append(Spacer(1, 0.3*cm))

    # Validator's signed opinion (if set)
    validation_opinion = (getattr(model, "validation_opinion", None) or "").strip()
    if validation_opinion:
        story.append(Paragraph("Validator's Opinion", S["h3"]))
        story.append(Table(
            [[Paragraph(validation_opinion,
                style("Opinion", fontName="Helvetica-Oblique", fontSize=9.5,
                      textColor=rgb(*NAVY), leading=14))]],
            colWidths=[page_w - 3.6*cm],
            style=TableStyle([
                ("BACKGROUND",    (0,0), (-1,-1), c_light),
                ("TOPPADDING",    (0,0), (-1,-1), 10),
                ("BOTTOMPADDING", (0,0), (-1,-1), 10),
                ("LEFTPADDING",   (0,0), (-1,-1), 12),
                ("RIGHTPADDING",  (0,0), (-1,-1), 12),
                ("BOX", (0,0), (-1,-1), 1, rgb(*TEAL)),
            ])
        ))
        story.append(Spacer(1, 0.4*cm))

    story.append(Paragraph("Recommended Next Steps", S["h3"]))
    critical_count = counts_by_sev.get("Critical", 0)
    high_count     = counts_by_sev.get("High", 0)

    if critical_count > 0:
        story.append(Paragraph(
            f"1. Address {critical_count} Critical finding(s) immediately — the model cannot be "
            f"deployed or relied upon for regulatory compliance until these are resolved.", S["body"]
        ))
    if high_count > 0:
        story.append(Paragraph(
            f"2. Resolve {high_count} High severity finding(s) before validation sign-off. "
            f"Each must be remediated or formally accepted with documented rationale.", S["body"]
        ))
    story.append(Paragraph(
        "3. Complete the regulatory compliance matrix — ensure all NOT MET requirements "
        "are addressed or escalated to the compliance function.", S["body"]
    ))
    story.append(Paragraph(
        "4. Schedule next review — models must be re-validated within 12 months or "
        "sooner if a material change occurs (FFFS 2017:11 §17).", S["body"]
    ))

    story.append(Spacer(1, 1*cm))
    story.append(Table(
        [[Paragraph("Validator sign-off", S["label"]), Paragraph("_" * 40, S["body"]),
          Paragraph("Date", S["label"]), Paragraph("_" * 20, S["body"])]],
        colWidths=[4*cm, 8*cm, 2*cm, 4*cm],
        style=TableStyle([
            ("TOPPADDING",    (0,0), (-1,-1), 10),
            ("BOTTOMPADDING", (0,0), (-1,-1), 10),
            ("LINEBELOW", (1,0), (1,0), 0.5, rgb(156,163,175)),
            ("LINEBELOW", (3,0), (3,0), 0.5, rgb(156,163,175)),
        ])
    ))

    # Build
    doc.build(story, onFirstPage=on_page, onLaterPages=on_page)
    return buffer.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# PPTX Report
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx_report(project, model, findings, compliance_results, compliance_summary_data) -> bytes:
    """
    Generate an executive PowerPoint presentation.
    Returns raw PPTX bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb_c(r, g, b):
        return RGBColor(r, g, b)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    W = prs.slide_width
    H = prs.slide_height

    def add_rect(slide, x, y, w, h, fill_rgb, radius=None):
        shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
        shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, font_name="Calibri", font_size=18,
                 bold=False, color=(255,255,255), align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(x, y, w, h)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return txb

    def inches(n):
        return Inches(n)

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)

    # Full navy background
    add_rect(slide, 0, 0, W, H, NAVY)

    # Teal accent bar on left
    add_rect(slide, 0, 0, Inches(0.12), H, TEAL)

    # Teal bar top-right corner
    add_rect(slide, W - Inches(3), 0, Inches(3), Inches(0.08), TEAL)

    # VALTIOR label
    add_text(slide, "VALTIOR", inches(0.4), inches(0.3), inches(5), inches(0.5),
             font_size=11, bold=True, color=(0,180,180))

    # Main title
    add_text(slide, "Model Validation Report", inches(0.4), inches(1.6), inches(8), inches(1.2),
             font_size=38, bold=True, color=WHITE)

    add_text(slide, project.name, inches(0.4), inches(2.9), inches(10), inches(0.7),
             font_size=20, color=(200, 220, 240))

    # Metadata box
    add_rect(slide, inches(0.4), inches(3.8), inches(5.5), inches(2.5), (20, 50, 80))

    meta_lines = [
        f"Institution:   {project.institution or '—'}",
        f"Model owner:  {project.owner or '—'}",
        f"Domain:        {project.domain}",
        f"Model version: v{model.version}",
        f"Date:          {datetime.now().strftime('%d %B %Y')}",
    ]
    add_text(slide, "\n".join(meta_lines), inches(0.7), inches(4.0), inches(5), inches(2.1),
             font_size=11, color=(200, 220, 240))

    # RAG badge
    rag = compliance_summary_data["rag"]
    rag_rgb = {"GREEN": (22,163,74), "AMBER": (234,88,12), "RED": (220,38,38)}[rag]
    add_rect(slide, inches(7.2), inches(3.8), inches(5.5), inches(1.2), rag_rgb)
    add_text(slide, f"Overall Status: {rag}", inches(7.4), inches(3.95), inches(5), inches(0.5),
             font_size=16, bold=True, color=WHITE)
    add_text(slide, compliance_summary_data["rag_label"], inches(7.4), inches(4.5), inches(5), inches(0.4),
             font_size=10, color=WHITE)

    add_text(slide, "CONFIDENTIAL", inches(0.4), inches(6.9), inches(12), inches(0.4),
             font_size=9, color=(100,130,160), align=PP_ALIGN.RIGHT)

    # ── Slide 2: Findings Overview ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, Inches(1.2), NAVY)
    add_text(slide, "Validation Overview", inches(0.4), inches(0.25), inches(10), inches(0.7),
             font_size=26, bold=True, color=WHITE)

    counts_by_sev = {s: 0 for s in ["Critical","High","Medium","Low","Observation"]}
    for f in findings:
        sev = f.get("severity") if isinstance(f,dict) else getattr(f,"severity","")
        if sev in counts_by_sev:
            counts_by_sev[sev] += 1

    # Severity stat boxes
    box_x = inches(0.4)
    box_w = inches(2.3)
    box_gap = inches(0.25)
    for sev, count in counts_by_sev.items():
        sc = SEV_COLOURS.get(sev, (107,114,128))
        add_rect(slide, box_x, inches(1.5), box_w, inches(1.8), sc)
        add_text(slide, str(count), box_x, inches(1.65), box_w, inches(0.9),
                 font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, sev, box_x, inches(2.6), box_w, inches(0.5),
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        box_x += box_w + box_gap

    # Compliance summary
    comp = compliance_summary_data["counts"]
    add_rect(slide, inches(0.4), inches(3.6), inches(12.5), inches(0.7), LIGHT)
    add_text(slide,
             f"Swedish Law Compliance:  ✓ {comp['MET']} Met    ~ {comp['PARTIAL']} Partial    ✗ {comp['NOT_MET']} Not Met    ? {comp['UNCLEAR']} Unclear   |   {compliance_summary_data['met_pct']}% overall",
             inches(0.6), inches(3.7), inches(12), inches(0.5),
             font_size=12, color=NAVY)

    # Top findings preview
    add_text(slide, "Critical & High Findings", inches(0.4), inches(4.5), inches(12), inches(0.4),
             font_size=13, bold=True, color=NAVY)

    critical_high = [
        f for f in findings
        if (f.get("severity") if isinstance(f,dict) else getattr(f,"severity","")) in ("Critical","High")
    ][:4]

    row_y = inches(5.0)
    for f in critical_high:
        sev   = f.get("severity") if isinstance(f,dict) else getattr(f,"severity","")
        title = f.get("title")    if isinstance(f,dict) else getattr(f,"title","")
        sc    = SEV_COLOURS.get(sev, (107,114,128))
        add_rect(slide, inches(0.4), row_y, inches(1.5), inches(0.35), sc)
        add_text(slide, sev, inches(0.4), row_y + inches(0.04), inches(1.5), inches(0.3),
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title[:90], inches(2.1), row_y, inches(11), inches(0.35),
                 font_size=10, color=NAVY)
        row_y += inches(0.42)

    # ── Slide 3: Regulatory Compliance Matrix ─────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, Inches(1.2), NAVY)
    add_text(slide, "Swedish Regulatory Compliance", inches(0.4), inches(0.25), inches(10), inches(0.7),
             font_size=26, bold=True, color=WHITE)
    add_text(slide, "PTL (2017:630) · FFFS 2017:11 · SIMPT Vägledning 2024",
             inches(0.4), inches(0.88), inches(10), inches(0.3),
             font_size=10, color=(0,180,180))

    row_y = inches(1.35)
    col_w = [inches(1.5), inches(3.5), inches(6.8), inches(1.5)]
    headers = ["ID", "Law", "Requirement", "Status"]
    for i, h in enumerate(headers):
        x = inches(0.4) + sum(col_w[:i])
        add_rect(slide, x, row_y, col_w[i], inches(0.38), TEAL)
        add_text(slide, h, x + inches(0.06), row_y + inches(0.06), col_w[i], inches(0.28),
                 font_size=9, bold=True, color=WHITE)

    row_y += inches(0.42)
    for idx, r in enumerate(compliance_results):
        bg = LIGHT if idx % 2 == 0 else WHITE
        add_rect(slide, inches(0.4), row_y, W - inches(0.8), inches(0.35), bg)
        vals = [r["id"], r["law"], r["short"], STATUS_LABELS[r["status"]]]
        for i, val in enumerate(vals):
            x = inches(0.4) + sum(col_w[:i])
            color = NAVY
            if i == 3:  # status column
                color = STATUS_COLOURS.get(r["status"], (107,114,128))
            add_text(slide, val, x + inches(0.06), row_y + inches(0.04), col_w[i], inches(0.28),
                     font_size=8.5, color=color,
                     bold=(i == 3))
        row_y += inches(0.38)

    # ── Slides 4+: Critical and High findings (2 per slide) ───────────────────
    sorted_findings = sorted(
        findings,
        key=lambda f: {"Critical":0,"High":1,"Medium":2,"Low":3,"Observation":4}.get(
            f.get("severity") if isinstance(f,dict) else getattr(f,"severity",""), 99
        )
    )

    chunked = [sorted_findings[i:i+2] for i in range(0, min(len(sorted_findings), 8), 2)]

    for chunk in chunked:
        slide = prs.slides.add_slide(blank_layout)
        add_rect(slide, 0, 0, W, Inches(1.0), NAVY)
        add_text(slide, "Validation Findings", inches(0.4), inches(0.2), inches(10), inches(0.6),
                 font_size=24, bold=True, color=WHITE)

        y_pos = inches(1.15)
        for f in chunk:
            sev   = f.get("severity")   if isinstance(f,dict) else getattr(f,"severity","")
            title = f.get("title")      if isinstance(f,dict) else getattr(f,"title","")
            desc  = f.get("description") if isinstance(f,dict) else getattr(f,"description","")
            rec   = f.get("recommendation") if isinstance(f,dict) else getattr(f,"recommendation","")
            reg   = f.get("regulatory_reference") if isinstance(f,dict) else getattr(f,"regulatory_reference","")
            sc    = SEV_COLOURS.get(sev, (107,114,128))

            add_rect(slide, inches(0.4), y_pos, W - inches(0.8), inches(0.45), sc)
            add_text(slide, f"{sev}  |  {title[:80]}", inches(0.6), y_pos + inches(0.07),
                     W - inches(1.2), inches(0.35), font_size=13, bold=True, color=WHITE)

            add_rect(slide, inches(0.4), y_pos + inches(0.45), W - inches(0.8), inches(0.28), LIGHT)
            add_text(slide, "Finding", inches(0.6), y_pos + inches(0.49),
                     inches(2), inches(0.22), font_size=8, bold=True, color=TEAL)

            add_text(slide, (desc or "")[:220], inches(0.6), y_pos + inches(0.75),
                     W - inches(1.2), inches(0.65), font_size=9.5, color=NAVY)

            add_text(slide, "Recommendation", inches(0.6), y_pos + inches(1.42),
                     inches(2.5), inches(0.22), font_size=8, bold=True, color=TEAL)
            add_text(slide, (rec or "")[:200], inches(0.6), y_pos + inches(1.65),
                     W - inches(1.2), inches(0.5), font_size=9.5, color=NAVY)

            if reg:
                add_text(slide, f"Ref: {reg}", inches(0.6), y_pos + inches(2.18),
                         W - inches(1.2), inches(0.28), font_size=8, color=(0,123,123))

            y_pos += inches(2.65)

    # ── Last Slide: Next Steps ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, NAVY)
    add_rect(slide, 0, 0, Inches(0.12), H, TEAL)

    add_text(slide, "Recommended Next Steps", inches(0.4), inches(0.4), inches(12), inches(0.8),
             font_size=30, bold=True, color=WHITE)

    steps = [
        ("1", "Resolve Critical findings", "The model cannot be used for regulatory compliance until Critical findings are addressed."),
        ("2", "Remediate High findings", "High severity gaps must be resolved or formally accepted before validation sign-off."),
        ("3", "Complete compliance matrix", "All NOT MET regulatory requirements need remediation plans with owners and deadlines."),
        ("4", "Schedule next review", "Set a next validation date (max 12 months) and define material change triggers."),
    ]

    step_y = inches(1.4)
    for num, title, desc in steps:
        add_rect(slide, inches(0.4), step_y, inches(0.7), inches(0.7), TEAL)
        add_text(slide, num, inches(0.4), step_y + inches(0.12), inches(0.7), inches(0.5),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, inches(1.3), step_y, inches(11), inches(0.4),
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, desc, inches(1.3), step_y + inches(0.42), inches(11), inches(0.35),
                 font_size=10, color=(180, 210, 240))
        step_y += inches(1.15)

    # Validator's opinion box (if set)
    validation_opinion_pptx = (getattr(model, "validation_opinion", None) or "").strip()
    if validation_opinion_pptx:
        add_rect(slide, inches(0.4), inches(6.0), W - inches(0.8), inches(0.75), TEAL)
        add_text(slide, "Validator's Opinion", inches(0.6), inches(6.04),
                 inches(3), inches(0.28), font_size=8, bold=True, color=WHITE)
        add_text(slide, validation_opinion_pptx[:250], inches(0.6), inches(6.3),
                 W - inches(1.2), inches(0.38), font_size=8.5, color=(220,240,240))

    add_text(slide, f"Generated by Valtior · {datetime.now().strftime('%d %B %Y')} · CONFIDENTIAL",
             inches(0.4), inches(6.9), inches(12), inches(0.4),
             font_size=9, color=(80,120,150), align=PP_ALIGN.RIGHT)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
